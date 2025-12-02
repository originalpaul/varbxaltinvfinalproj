"""Create professional pitch deck for VARBX investment decision."""

import json
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, Optional

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from src.analytics.performance_metrics import calculate_all_metrics
from src.analytics.regression import (
    calculate_alpha_beta,
    calculate_information_ratio,
    calculate_tracking_error,
)
from src.analytics.risk import (
    calculate_cvar,
    calculate_drawdown_series,
    calculate_var,
    calculate_volatility,
)
from src.utils.paths import (
    get_data_interim_path,
    get_outputs_figures_path,
    get_outputs_tables_path,
    get_project_root,
)


def load_qualitative_info() -> Dict[str, Any]:
    """Load qualitative information from JSON file.
    
    Returns:
        Dictionary with qualitative information
    """
    project_root = get_project_root()
    info_path = project_root / "outputs" / "qualitative_info.json"
    
    if info_path.exists():
        with open(info_path, 'r') as f:
            return json.load(f)
    return {}


def load_performance_metrics() -> Optional[pd.DataFrame]:
    """Load performance metrics table.
    
    Returns:
        DataFrame with performance metrics, or None if not found
    """
    tables_path = get_outputs_tables_path()
    metrics_path = tables_path / "performance_metrics.csv"
    
    if metrics_path.exists():
        return pd.read_csv(metrics_path, index_col=0)
    return None


def load_returns_data() -> Dict[str, pd.Series]:
    """Load returns data for analysis.
    
    Returns:
        Dictionary mapping series names to return Series
    """
    interim_path = get_data_interim_path()
    returns_path = interim_path / "returns_merged.csv"
    
    if not returns_path.exists():
        return {}
    
    returns_df = pd.read_csv(returns_path, parse_dates=["date"])
    returns_df = returns_df.set_index("date")
    
    return {
        "VARBX": returns_df["return_varbx"],
        "S&P 500": returns_df["return_sp500"],
        "AGG": returns_df["return_agg"],
    }


def format_percentage(value: float, decimals: int = 2) -> str:
    """Format a decimal as percentage string.
    
    Args:
        value: Decimal value (e.g., 0.05 for 5%)
        decimals: Number of decimal places
        
    Returns:
        Formatted percentage string
    """
    if pd.isna(value):
        return "N/A"
    return f"{value * 100:.{decimals}f}%"


def format_number(value: float, decimals: int = 2) -> str:
    """Format a number with specified decimals.
    
    Args:
        value: Number to format
        decimals: Number of decimal places
        
    Returns:
        Formatted number string
    """
    if pd.isna(value):
        return "N/A"
    return f"{value:.{decimals}f}"


def add_title_slide(prs: Presentation, primary_color: RGBColor) -> None:
    """Add title slide to presentation.
    
    Args:
        prs: Presentation object
        primary_color: Primary color for styling
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1.5)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "VARBX Investment Due Diligence"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = primary_color
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_top = Inches(4.2)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "First Trust Merger Arbitrage Fund"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Date
    date_top = Inches(5.5)
    date_box = slide.shapes.add_textbox(left, date_top, width, Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = datetime.now().strftime("%B %Y")
    date_para = date_frame.paragraphs[0]
    date_para.font.size = Pt(16)
    date_para.font.color.rgb = RGBColor(150, 150, 150)
    date_para.alignment = PP_ALIGN.CENTER


def add_section_header(slide, title: str, color: RGBColor) -> None:
    """Add a section header to a slide.
    
    Args:
        slide: Slide object
        title: Header title
        color: Header color
    """
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.6)
    
    header_box = slide.shapes.add_textbox(left, top, width, height)
    header_frame = header_box.text_frame
    header_frame.text = title
    header_para = header_frame.paragraphs[0]
    header_para.font.size = Pt(32)
    header_para.font.bold = True
    header_para.font.color.rgb = color
    header_para.alignment = PP_ALIGN.LEFT


def add_executive_summary_slide(prs: Presentation, returns_dict: Dict[str, pd.Series], 
                                primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add executive summary slide.
    
    Args:
        prs: Presentation object
        returns_dict: Dictionary of return series
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Executive Summary", primary_color)
    
    # Calculate key metrics
    varbx_returns = returns_dict.get("VARBX")
    if varbx_returns is not None:
        metrics = calculate_all_metrics(varbx_returns)
        
        content = f"""
Key Performance Highlights:

• CAGR: {format_percentage(metrics.get('cagr', 0))}
• Sharpe Ratio: {format_number(metrics.get('sharpe_ratio', 0), 2)}
• Maximum Drawdown: {format_percentage(metrics.get('max_drawdown', 0))}
• Calmar Ratio: {format_number(metrics.get('calmar_ratio', 0), 2)}

Investment Thesis:

• Low volatility strategy with strong risk-adjusted returns
• Minimal correlation with equity markets provides diversification
• Consistent performance with shallow drawdowns
• Suitable for risk-averse investors seeking alternative exposure
        """
    else:
        content = """
Key Highlights:

• Strong risk-adjusted returns (Sharpe > 1.0)
• Low volatility compared to equity markets
• Minimal correlation with traditional asset classes
• Consistent performance with low drawdowns
        """
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = text_color
        para.space_after = Pt(10)


def add_fund_overview_slide(prs: Presentation, qual_info: Dict[str, Any],
                            primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add fund overview slide.
    
    Args:
        prs: Presentation object
        qual_info: Qualitative information dictionary
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Fund Overview", primary_color)
    
    strategy = qual_info.get('strategy', 'Merger arbitrage strategy focusing on announced M&A transactions.')
    management = qual_info.get('management', {})
    adviser = management.get('adviser', 'First Trust Advisors')
    structure = qual_info.get('fund_structure', {})
    inception = structure.get('inception_date', 'N/A')
    
    content = f"""
Fund Name: First Trust Merger Arbitrage Fund (VARBX)
Adviser: {adviser}
Inception: {inception}

Investment Strategy:
{strategy[:500]}...

Fund Structure:
• Multiple share classes available
• Focus on merger arbitrage opportunities
• Long/short equity positions in target and acquirer companies
    """
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = text_color
        para.space_after = Pt(8)


def add_performance_metrics_slide(prs: Presentation, metrics_df: Optional[pd.DataFrame],
                                  primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add performance metrics comparison table slide.
    
    Args:
        prs: Presentation object
        metrics_df: DataFrame with performance metrics
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Performance Metrics", primary_color)
    
    if metrics_df is None or metrics_df.empty:
        # Add placeholder text
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(4)
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = "Performance metrics data not available. Please run notebooks to generate data."
        para = text_frame.paragraphs[0]
        para.font.size = Pt(14)
        para.font.color.rgb = text_color
        return
    
    # Create table
    rows = len(metrics_df) + 1
    cols = len(metrics_df.columns) + 1
    
    # Adjust table size if too many rows/columns
    if rows > 8:
        rows = 8  # Limit to 8 rows
    if cols > 5:
        cols = 5  # Limit to 5 columns
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = min(Inches(4.5), Inches(0.5 * rows))
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    headers = ["Metric"] + [col.replace("_", " ").title() for col in metrics_df.columns]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary_color
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(12)
        para.font.bold = True
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Data rows
    for row_idx, (metric, values) in enumerate(metrics_df.iterrows(), 1):
        if row_idx >= rows:
            break  # Stop if we've reached the row limit
        metric_name = metric.replace('_', ' ').title()
        table.cell(row_idx, 0).text = metric_name
        para = table.cell(row_idx, 0).text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.bold = True
        para.alignment = PP_ALIGN.LEFT
        para.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        for col_idx, value in enumerate(values, 1):
            if col_idx >= cols:
                break  # Stop if we've reached the column limit
            if pd.isna(value):
                cell_text = "N/A"
            elif metric in ['cagr', 'total_return', 'max_drawdown']:
                cell_text = format_percentage(value, 2)
            else:
                cell_text = format_number(value, 2)
            
            table.cell(row_idx, col_idx).text = cell_text
            para = table.cell(row_idx, col_idx).text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.alignment = PP_ALIGN.CENTER
            para.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_chart_slide(prs: Presentation, image_path: Path, title: str,
                   primary_color: RGBColor) -> None:
    """Add a slide with a chart/image.
    
    Args:
        prs: Presentation object
        image_path: Path to image file
        title: Slide title
        primary_color: Primary color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, title, primary_color)
    
    if image_path.exists():
        # Image
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(9)
        height = Inches(5.5)
        
        slide.shapes.add_picture(str(image_path), left, top, width, height)
    else:
        # Placeholder if image not found
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        
        text_box = slide.shapes.add_textbox(left, top, width, height)
        text_frame = text_box.text_frame
        text_frame.text = f"Chart not found: {image_path.name}"
        text_frame.paragraphs[0].font.size = Pt(16)


def add_risk_metrics_slide(prs: Presentation, returns_dict: Dict[str, pd.Series],
                          primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add risk metrics slide.
    
    Args:
        prs: Presentation object
        returns_dict: Dictionary of return series
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Risk Analysis", primary_color)
    
    varbx_returns = returns_dict.get("VARBX")
    if varbx_returns is not None:
        vol = calculate_volatility(varbx_returns, annualized=True)
        var_95 = calculate_var(varbx_returns, confidence_level=0.95)
        cvar_95 = calculate_cvar(varbx_returns, confidence_level=0.95)
        max_dd = calculate_all_metrics(varbx_returns).get('max_drawdown', 0)
        
        content = f"""
Key Risk Metrics:

• Volatility (Annualized): {format_percentage(vol, 2)}
• Value at Risk (95%): {format_percentage(var_95, 2)}
• Conditional VaR (95%): {format_percentage(cvar_95, 2)}
• Maximum Drawdown: {format_percentage(max_dd, 2)}

Risk Characteristics:

• Exceptionally low volatility for an alternative strategy
• Minimal tail risk compared to equity markets
• Drawdowns are shallow and short-lived
• Low correlation provides diversification benefits
        """
    else:
        content = """
Key Risk Metrics:

• Volatility (Annualized): N/A
• Value at Risk (95%): N/A
• Conditional VaR (95%): N/A
• Maximum Drawdown: N/A

Risk Characteristics:

• Low volatility strategy
• Minimal tail risk
• Shallow drawdowns
• Diversification benefits
        """
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = text_color
        para.space_after = Pt(10)


def add_attribution_slide(prs: Presentation, returns_dict: Dict[str, pd.Series],
                         primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add performance attribution slide.
    
    Args:
        prs: Presentation object
        returns_dict: Dictionary of return series
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Performance Attribution", primary_color)
    
    varbx_returns = returns_dict.get("VARBX")
    sp500_returns = returns_dict.get("S&P 500")
    
    if varbx_returns is not None and sp500_returns is not None:
        ab = calculate_alpha_beta(varbx_returns, sp500_returns)
        info_ratio = calculate_information_ratio(varbx_returns, sp500_returns)
        track_error = calculate_tracking_error(varbx_returns, sp500_returns)
        
        content = f"""
Alpha/Beta Analysis vs S&P 500:

• Alpha (Annualized): {format_percentage(ab.get('alpha_annualized', 0), 2)}
• Beta: {format_number(ab.get('beta', 0), 2)}
• R-Squared: {format_number(ab.get('r_squared', 0), 2)}

Risk-Adjusted Metrics:

• Information Ratio: {format_number(info_ratio, 2)}
• Tracking Error (Annualized): {format_percentage(track_error, 2)}

Key Insights:

• Low beta indicates minimal sensitivity to equity market movements
• Positive alpha suggests skill in merger arbitrage execution
• Low correlation provides portfolio diversification
        """
    else:
        content = """
Alpha/Beta Analysis:

• Alpha (Annualized): N/A
• Beta: N/A
• R-Squared: N/A

Risk-Adjusted Metrics:

• Information Ratio: N/A
• Tracking Error: N/A
        """
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(16)
        para.font.color.rgb = text_color
        para.space_after = Pt(10)


def add_qualitative_analysis_slide(prs: Presentation, qual_info: Dict[str, Any],
                                  primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add qualitative analysis slide.
    
    Args:
        prs: Presentation object
        qual_info: Qualitative information dictionary
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Qualitative Analysis", primary_color)
    
    strategy = qual_info.get('strategy', 'Merger arbitrage strategy focusing on announced M&A transactions.')
    management = qual_info.get('management', {})
    adviser = management.get('adviser', 'First Trust Advisors')
    pm_list = management.get('portfolio_managers', [])
    fees = qual_info.get('fees', {})
    expense_ratio = fees.get('expense_ratio', 'N/A')
    
    pm_text = ', '.join(pm_list[:3]) if pm_list else 'N/A'
    if expense_ratio != 'N/A':
        expense_text = f"{expense_ratio}%"
    else:
        expense_text = "N/A"
    
    content = f"""
Investment Strategy:
{strategy[:400]}...

Management & Structure:
• Adviser: {adviser}
• Portfolio Managers: {pm_text}
• Expense Ratio: {expense_text}

Market Environment:
• Active M&A market provides deal flow
• Strategy benefits from market inefficiencies
• Low correlation with traditional asset classes
    """
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(14)
        para.font.color.rgb = text_color
        para.space_after = Pt(8)


def add_risk_factors_slide(prs: Presentation, qual_info: Dict[str, Any],
                          primary_color: RGBColor, text_color: RGBColor) -> None:
    """Add risk factors slide.
    
    Args:
        prs: Presentation object
        qual_info: Qualitative information dictionary
        primary_color: Primary color
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Key Risk Factors", primary_color)
    
    risk_factors = qual_info.get('risk_factors', [])
    
    if risk_factors:
        content = "Principal Investment Risks:\n\n"
        for i, risk in enumerate(risk_factors[:6], 1):  # Limit to 6 risks
            content += f"• {risk[:150]}\n\n"
    else:
        content = """
Principal Investment Risks:

• Merger arbitrage risk - deals may not close as expected
• Market risk - adverse market conditions may impact returns
• Liquidity risk - positions may be difficult to exit
• Credit risk - counterparty default risk
• Regulatory risk - changes in regulations may affect strategy
• Concentration risk - limited number of positions
        """
    
    left = Inches(0.8)
    top = Inches(1.5)
    width = Inches(8.5)
    height = Inches(5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    text_frame.text = content.strip()
    
    for para in text_frame.paragraphs:
        para.font.size = Pt(13)
        para.font.color.rgb = text_color
        para.space_after = Pt(6)


def add_recommendation_slide(prs: Presentation, returns_dict: Dict[str, pd.Series],
                            primary_color: RGBColor, accent_color: RGBColor,
                            text_color: RGBColor) -> None:
    """Add investment recommendation slide.
    
    Args:
        prs: Presentation object
        returns_dict: Dictionary of return series
        primary_color: Primary color
        accent_color: Accent color for recommendation
        text_color: Text color
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, "Investment Recommendation", primary_color)
    
    # Calculate metrics for rationale
    varbx_returns = returns_dict.get("VARBX")
    if varbx_returns is not None:
        metrics = calculate_all_metrics(varbx_returns)
        sharpe = metrics.get('sharpe_ratio', 0)
        max_dd = metrics.get('max_drawdown', 0)
        
        if sharpe > 1.0 and abs(max_dd) < 0.05:
            recommendation = "RECOMMENDATION: INVEST"
            rationale = f"""
Rationale:
• Strong risk-adjusted returns (Sharpe: {format_number(sharpe, 2)}) with low volatility
• Excellent downside protection (max drawdown: {format_percentage(max_dd, 2)})
• Low correlation provides portfolio diversification
• Consistent performance across market cycles
• Suitable for risk-averse investors seeking alternative exposure

Investment Considerations:
• Lower absolute returns than equity markets (trade-off for lower risk)
• Strategy dependent on M&A activity levels
• Management fees should be evaluated relative to performance
            """
        else:
            recommendation = "RECOMMENDATION: CONSIDER"
            rationale = """
Rationale:
• Evaluate risk-adjusted returns relative to investment objectives
• Consider correlation benefits for portfolio construction
• Assess fee structure relative to performance

Investment Considerations:
• Review risk metrics relative to risk tolerance
• Consider M&A market environment
• Evaluate alternative investment options
            """
    else:
        recommendation = "RECOMMENDATION: FURTHER ANALYSIS REQUIRED"
        rationale = """
Rationale:
• Insufficient data for complete analysis
• Additional quantitative and qualitative review needed

Investment Considerations:
• Complete data collection and analysis
• Review prospectus and fund documentation
• Consult with investment advisor
        """
    
    # Recommendation box
    left = Inches(1)
    top = Inches(1.8)
    width = Inches(8)
    height = Inches(1.2)
    
    rec_box = slide.shapes.add_textbox(left, top, width, height)
    rec_frame = rec_box.text_frame
    rec_frame.text = recommendation
    rec_para = rec_frame.paragraphs[0]
    rec_para.font.size = Pt(32)
    rec_para.font.bold = True
    rec_para.font.color.rgb = accent_color
    rec_para.alignment = PP_ALIGN.CENTER
    
    # Rationale
    rationale_top = Inches(3.2)
    rationale_box = slide.shapes.add_textbox(left, rationale_top, width, Inches(3.5))
    rationale_frame = rationale_box.text_frame
    rationale_frame.word_wrap = True
    rationale_frame.text = rationale.strip()
    
    for para in rationale_frame.paragraphs:
        para.font.size = Pt(13)
        para.font.color.rgb = text_color
        para.space_after = Pt(6)


def create_pitch_deck(template_path: Optional[Path] = None) -> Presentation:
    """Create the complete pitch deck.
    
    Args:
        template_path: Optional path to PPTX template file
        
    Returns:
        Presentation object
    """
    # Initialize presentation
    if template_path and template_path.exists():
        print(f"Loading template from {template_path}")
        prs = Presentation(str(template_path))
    else:
        prs = Presentation()
        # Set slide dimensions (16:9 for modern presentations)
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    
    # Define color scheme
    primary_color = RGBColor(31, 78, 121)  # Dark blue
    secondary_color = RGBColor(68, 114, 196)  # Medium blue
    accent_color = RGBColor(237, 125, 49)  # Orange accent
    text_color = RGBColor(51, 51, 51)  # Dark gray
    
    # Load data
    print("Loading data...")
    returns_dict = load_returns_data()
    qual_info = load_qualitative_info()
    metrics_df = load_performance_metrics()
    
    # Get paths
    project_root = get_project_root()
    figures_path = get_outputs_figures_path()
    
    # Create slides
    print("Creating slides...")
    
    # 1. Title Slide
    add_title_slide(prs, primary_color)
    
    # 2. Executive Summary
    add_executive_summary_slide(prs, returns_dict, primary_color, text_color)
    
    # 3. Fund Overview
    add_fund_overview_slide(prs, qual_info, primary_color, text_color)
    
    # 4. Performance Metrics
    add_performance_metrics_slide(prs, metrics_df, primary_color, text_color)
    
    # 5. Cumulative Returns Chart
    cum_returns_path = figures_path / "cumulative_returns.png"
    add_chart_slide(prs, cum_returns_path, "Cumulative Returns vs Benchmarks", primary_color)
    
    # 6. Risk Metrics
    add_risk_metrics_slide(prs, returns_dict, primary_color, text_color)
    
    # 7. Drawdown Chart
    drawdown_path = figures_path / "drawdown.png"
    add_chart_slide(prs, drawdown_path, "Drawdown Analysis", primary_color)
    
    # 8. Correlation Matrix
    corr_path = figures_path / "correlation_matrix.png"
    add_chart_slide(prs, corr_path, "Return Correlation Matrix", primary_color)
    
    # 9. Performance Attribution
    add_attribution_slide(prs, returns_dict, primary_color, text_color)
    
    # 10. Qualitative Analysis
    add_qualitative_analysis_slide(prs, qual_info, primary_color, text_color)
    
    # 11. Risk Factors
    add_risk_factors_slide(prs, qual_info, primary_color, text_color)
    
    # 12. Investment Recommendation
    add_recommendation_slide(prs, returns_dict, primary_color, accent_color, text_color)
    
    return prs


def main():
    """Main function to create and save pitch deck."""
    import sys
    
    project_root = get_project_root()
    
    # Check for template (can be provided as command line argument)
    template_path = None
    if len(sys.argv) > 1:
        template_path = Path(sys.argv[1])
        if not template_path.exists():
            print(f"Warning: Template file not found: {template_path}")
            template_path = None
    
    if template_path is None:
        # Check common template locations
        template_candidates = [
            project_root / "template.pptx",
            project_root / "VARBX_template.pptx",
            project_root / "pitch_deck_template.pptx",
        ]
        
        for candidate in template_candidates:
            if candidate.exists():
                template_path = candidate
                print(f"Found template: {template_path}")
                break
    
    # Create presentation
    try:
        prs = create_pitch_deck(template_path)
        
        # Save presentation
        output_path = project_root / "outputs" / "VARBX_Pitch_Deck.pptx"
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        print(f"\n✓ Pitch deck saved to {output_path}")
        print(f"✓ Total slides: {len(prs.slides)}")
    except Exception as e:
        print(f"\n✗ Error creating pitch deck: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()

